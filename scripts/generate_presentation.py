"""Generate the updated SEC EDGAR Analyst Agent Capstone presentation PowerPoint file (.pptx).

Uses python-pptx to build a 21-slide widescreen presentation matching Google Cloud dark theme aesthetics,
complete with technical trade-off matrices, backable metrics, and verbatim speaker notes on every slide.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- COLOR PALETTE (Google Cloud Dark Theme) ---
COLOR_BG = RGBColor(15, 23, 42)         # #0F172A Dark Slate
COLOR_CARD_BG = RGBColor(30, 41, 59)     # #1E293B Card Fill
COLOR_CARD_BORDER = RGBColor(51, 65, 85) # #334155 Card Border
COLOR_TEXT_PRIMARY = RGBColor(255, 255, 255) # Pure White
COLOR_TEXT_MUTED = RGBColor(148, 163, 184)   # #94A3B8 Slate Gray
COLOR_GCP_BLUE = RGBColor(66, 133, 244)   # #4285F4 Google Blue
COLOR_GCP_GREEN = RGBColor(52, 168, 83)   # #34A853 Google Green
COLOR_GCP_YELLOW = RGBColor(251, 188, 4)  # #FBBC04 Google Yellow
COLOR_GCP_RED = RGBColor(234, 67, 53)     # #EA4335 Google Red
COLOR_CARD_GREEN_BG = RGBColor(20, 50, 30) # Soft Dark Green Box
COLOR_CARD_RED_BG = RGBColor(50, 20, 20)   # Soft Dark Red Box


def create_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_header(slide, title: str, subtitle: str, category: str = "Confidential + Proprietary | Google Cloud"):
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.2))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p0 = tf.paragraphs[0]
        p0.text = category.upper()
        p0.font.size = Pt(10)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_GCP_BLUE
        p0.space_after = Pt(4)

        p1 = tf.add_paragraph()
        p1.text = title
        p1.font.size = Pt(24)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_PRIMARY
        p1.space_after = Pt(2)

        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(13)
        p2.font.color.rgb = COLOR_GCP_BLUE

    def add_footer(slide):
        bar_y = Inches(7.35)
        w_segment = Inches(13.333 / 4)
        colors = [COLOR_GCP_BLUE, COLOR_GCP_RED, COLOR_GCP_YELLOW, COLOR_GCP_GREEN]
        for i, col in enumerate(colors):
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, w_segment * i, bar_y, w_segment, Inches(0.15))
            shape.fill.solid()
            shape.fill.fore_color.rgb = col
            shape.line.fill.background()

        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(11.733), Inches(0.3))
        tf = tx_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Confidential + Proprietary | Google Cloud Forward Deployed Engineering"
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_TEXT_MUTED

    def add_card(slide, left, top, width, height, title="", bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)

        if title:
            tx_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.2), Inches(width - 0.4), Inches(0.5))
            tf = tx_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = COLOR_TEXT_PRIMARY
        return shape

    def add_bullet_points(slide, left, top, width, height, bullets, font_size=13):
        tx_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        for i, text in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {text}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = COLOR_TEXT_PRIMARY
            p.space_after = Pt(8)

    def set_speaker_notes(slide, notes_text: str):
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes_text

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1)

    tx_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf = tx_box.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "CONFIDENTIAL + PROPRIETARY | GOOGLE CLOUD GTM FORWARD DEPLOYED ENGINEERING"
    p0.font.size = Pt(11)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_TEXT_MUTED
    p0.space_after = Pt(20)

    p1 = tf.add_paragraph()
    p1.text = "SEC EDGAR Analyst Agent"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT_PRIMARY
    p1.space_after = Pt(10)

    p2 = tf.add_paragraph()
    p2.text = "Grounded Financial Intelligence & Automated Reasoning Engine on Google Cloud"
    p2.font.size = Pt(22)
    p2.font.color.rgb = COLOR_GCP_BLUE
    p2.space_after = Pt(40)

    p3 = tf.add_paragraph()
    p3.text = "Charles Wang  |  Forward Deployed Engineer"
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_TEXT_MUTED

    add_footer(slide1)
    set_speaker_notes(slide1, 
        "Welcome everyone. Today I'm presenting the SEC EDGAR Analyst Agent, an enterprise-grade financial reasoning engine built on Google Cloud Vertex AI and the Agent Development Kit (ADK). This project transforms manual SEC filing analysis into a 100% mathematically precise, grounded agentic solution."
    )

    # ==========================================
    # SLIDE 2: Introduction
    # ==========================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2)
    add_header(slide2, "Introduction", "Engineering Origins & Background")

    add_card(slide2, 0.8, 1.8, 3.6, 4.8, "Charles Wang")
    add_bullet_points(slide2, 1.0, 2.5, 3.2, 3.8, [
        "Forward Deployed Engineer at Google Cloud",
        "Prior: FDE at Palantir & Strategic Projects Lead at Scale AI",
        "Focused on Enterprise AI, RAG Systems, & Agent Architecture"
    ], font_size=12)

    add_card(slide2, 4.7, 1.8, 3.8, 4.8, "Education & Technical Focus")
    add_bullet_points(slide2, 4.9, 2.5, 3.4, 3.8, [
        "B.S. in Computer Science Engineering — University of Michigan",
        "MBA — University of Chicago Booth School of Business",
        "Specialization: Financial Modeling & AI System Design"
    ], font_size=12)

    add_card(slide2, 8.8, 1.8, 3.7, 4.8, "Capstone Scope")
    add_bullet_points(slide2, 9.0, 2.5, 3.3, 3.8, [
        "Domain: Financial Services Industry (FSI)",
        "Goal: Automate multi-statement SEC 10-K extraction & variance analysis",
        "Focus: Zero-hallucination math & auditability"
    ], font_size=12)

    add_footer(slide2)
    set_speaker_notes(slide2,
        "Briefly on my background: I'm a Forward Deployed Engineer here at Google Cloud, with prior experience at Palantir and Scale AI. I hold a Computer Science degree from Michigan and an MBA from Chicago Booth. Let's dive straight into the business problem."
    )

    # ==========================================
    # SLIDE 3: Agenda & Story Map
    # ==========================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3)
    add_header(slide3, "Agenda & Story Map", "Roadmap of Presentation Sections")

    sections = [
        ("01", "Business Problem & User Journey", "Manual 10-K pain & workflow transformation"),
        ("02", "Live Product Demo & Solution Overview", "Split-pane UI & 4 Technical Pillars"),
        ("03", "Architecture Deep Dives & Trade-Offs", "ADK Supervisor, Hybrid RAG, & Decoupled Math"),
        ("04", "Expanded Evaluation Framework", "Math Assertions, LLM-as-a-Judge, & ADK Eval Harness"),
        ("05", "TCO, Scalability & Strategic ROI", "Context Caching (CAG), BigQuery Telemetry, & Net Value")
    ]

    top_pos = 1.8
    for num, title, desc in sections:
        add_card(slide3, 0.8, top_pos, 11.733, 0.85)
        tx_box = slide3.shapes.add_textbox(Inches(1.0), Inches(top_pos + 0.15), Inches(11.3), Inches(0.55))
        tf = tx_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{num}   |   {title}  —  {desc}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_PRIMARY
        top_pos += 1.0

    add_footer(slide3)
    set_speaker_notes(slide3,
        "Here is our agenda for today. We'll start with the baseline manual problem in financial analysis, walk through a live user journey and demo, dive deep into each architectural layer with explicit trade-offs, review our expanded evaluation harness, and conclude with TCO and ROI."
    )

    # ==========================================
    # SLIDE 4: Business Problem & Analyst Bottleneck
    # ==========================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4)
    add_header(slide4, "Baseline Business Problem — The Analyst Bottleneck", "The Manual Challenge of SEC Filing Parsing in Financial Services")

    add_card(slide4, 0.8, 1.8, 5.7, 4.8, "Current Manual Process (Legacy State)")
    add_bullet_points(slide4, 1.0, 2.5, 5.3, 3.8, [
        "2.5 Hours per Report: 150 minutes spent manually reading 300+ page 10-K filings and calculating margins.",
        "3,000 Hours Wasted Annually: Across 15 Senior Analysts producing ~80 comprehensive reports per year.",
        "$270,000 Direct Labor Spend: 3,000 hours at $90/hr loaded FSI senior analyst rate.",
        "High Risk of Arithmetic Errors: Manual spreadsheet copy-pasting creates material risk in valuation models.",
        "Decision Friction: Slow 3-day turnaround delays portfolio rebalancing and executive investment decisions."
    ], font_size=12)

    add_card(slide4, 6.8, 1.8, 5.7, 4.8, "Why Standard LLMs Fail in Finance")
    add_bullet_points(slide4, 7.0, 2.5, 5.3, 3.8, [
        "Probabilistic Math Hallucinations: Standard LLMs calculate math probabilistically, producing incorrect variance numbers.",
        "Lack of Paragraph-Level Auditability: Black-box text generation prevents verification back to source filing excerpts.",
        "Security & Data Privacy Risks: Public API LLM usage exposes proprietary analysis and violates compliance perimeters.",
        "High Pre-fill Latency: Processing massive 300-page 10-Ks leads to 15s+ TTFT without context caching."
    ], font_size=12)

    add_footer(slide4)
    set_speaker_notes(slide4,
        "Before discussing AI metrics, let's examine the baseline business reality. Senior financial analysts spend 2.5 hours reading 300-page SEC 10-Ks and manually copy-pasting numbers into spreadsheets. This wastes 3,000 hours and $270,000 annually per team. Worse, standard LLMs cannot solve this off-the-shelf because probabilistic math introduces unacceptable hallucinations."
    )

    # ==========================================
    # SLIDE 5: Target User Journey & Transformation
    # ==========================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5)
    add_header(slide5, "Target User Journey & Operational Transformation", "Reclaiming 2,900 Hours of Analyst Time via Grounded AI")

    add_card(slide5, 0.8, 1.8, 5.7, 4.8, "Legacy Manual Analyst Workflow")
    add_bullet_points(slide5, 1.0, 2.5, 5.3, 3.8, [
        "1. SEC EDGAR Search: Manually find and download multi-year 10-K PDFs.",
        "2. Document Parsing: Skim 300+ pages of footnotes, Item 7 MD&A, and financial statements.",
        "3. Spreadsheet Math: Copy numbers into Excel to calculate YoY variance and growth rates.",
        "4. Memo Drafting: Write summary reports manually in Word.",
        "Total Time: 150 minutes per report | High Error Risk"
    ], font_size=12)

    add_card(slide5, 6.8, 1.8, 5.7, 4.8, "Grounded Agentic Workflow (After)", border_color=COLOR_GCP_BLUE)
    add_bullet_points(slide5, 7.0, 2.5, 5.3, 3.8, [
        "1. Natural Language Prompt: 'Compare Tesla 2022 vs 2023 Revenue & MD&A Risks.'",
        "2. Autonomous ADK Orchestration: SQL metrics lookup + sandboxed Python math + Vertex AI search.",
        "3. Split-Pane Verification: Click inline citations to inspect exact 10-K source paragraphs.",
        "4. Human-in-the-Loop Export: One-click approval to generate verified report.",
        "Total Time: < 5 minutes per report | 100% Math Precision"
    ], font_size=12)

    add_footer(slide5)
    set_speaker_notes(slide5,
        "Here is the operational transformation. In the legacy state, an analyst manually opens EDGAR, copies metrics into Excel, and writes memos over 2.5 hours. In our grounded agentic state, the analyst types a natural language query, receives deterministic math and grounded citations in under 5 minutes, verifies source excerpts split-pane, and approves export."
    )

    # ==========================================
    # SLIDE 6: System Requirements & Production SLAs
    # ==========================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide6)
    add_header(slide6, "System Requirements & Production SLAs", "Core System Capabilities & Empirical Performance Thresholds")

    sla_cards = [
        ("100%", "Math Accuracy", "Zero-LLM sandboxed Python calculation engine"),
        ("≥ 92%", "Grounded Precision", "Paragraph-level citations to SEC 10-K sections"),
        ("< 3.0s", "Latency SLA", "1.31s measured benchmark execution average"),
        ("100%", "Compliance & Audit", "Model Armor guardrails & VPC-SC security perimeter")
    ]

    for i, (val, title, desc) in enumerate(sla_cards):
        left = 0.8 + (i * 3.0)
        add_card(slide6, left, 1.8, 2.7, 2.4)
        tx = slide6.shapes.add_textbox(Inches(left + 0.1), Inches(1.9), Inches(2.5), Inches(2.2))
        tf = tx.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.text = val
        p0.font.size = Pt(36)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_GCP_BLUE
        
        p1 = tf.add_paragraph()
        p1.text = title
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_PRIMARY

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    add_card(slide6, 0.8, 4.5, 11.733, 2.2, "Key Functional Requirements")
    add_bullet_points(slide6, 1.0, 5.0, 11.3, 1.5, [
        "Deterministic numerical calculations for verifiable financial reporting.",
        "Dual-path hybrid search merging BigQuery SQL metrics with Vertex AI Search vector chunks.",
        "Split-pane UI rendering grounded inline citations linking back to original GCS filings.",
        "Human-in-the-loop (HITL) approval gate for all report exports."
    ], font_size=12)

    add_footer(slide6)
    set_speaker_notes(slide6,
        "Our system is engineered against four non-negotiable SLAs: 100% math accuracy, 92%+ grounded precision, sub-3.0s latency, and 100% security compliance. We enforce these using sandboxed Python math, Vertex AI Hybrid Search, and Model Armor guardrails."
    )

    # ==========================================
    # SLIDE 7: Live Product Demo & Interface Experience
    # ==========================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide7)
    add_header(slide7, "Live Product Demo & Interface Experience", "Interactive Split-Pane View, Grounded Citations, & A2UI Protocol")

    add_card(slide7, 0.8, 1.8, 6.0, 4.8, "Analyst User Experience (UI Breakdown)")
    add_bullet_points(slide7, 1.0, 2.5, 5.6, 3.8, [
        "1. Split-Pane Architecture: Left pane for conversational AI narrative; right pane for full-text 10-K document inspection.",
        "2. Grounded Citation Links: Clicking highlighted text (<mark>) immediately jumps the right pane to the exact paragraph in GCS.",
        "3. Dynamic A2UI Protocol: Auto-renders structured JSON specs as dynamic formatted variance tables and bar charts.",
        "4. Human-in-the-Loop Gate: Mandatory confirmation modal before exporting financial reports to Cloud Storage buckets."
    ], font_size=12)

    add_card(slide7, 7.1, 1.8, 5.4, 4.8, "UI Architecture Mockup", border_color=COLOR_GCP_BLUE)
    tx = slide7.shapes.add_textbox(Inches(7.3), Inches(2.5), Inches(5.0), Inches(3.8))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "[ LEFT PANE: Agent Chat & A2UI ]\n• User Query\n• Grounded Narrative + Inline Citations\n• Dynamic Financial Variance Table\n\n[ RIGHT PANE: Grounded Document ]\n• Highlighted 10-K Excerpt (Item 7 MD&A)\n• Source: gs://sec-analyst-sec-reports/..."
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_GCP_BLUE

    add_footer(slide7)
    set_speaker_notes(slide7,
        "Here is the product interface. The analyst asks a question, and the agent responds with a grounded narrative and a dynamic A2UI variance table. When the analyst clicks any inline citation, the right-hand split-pane immediately navigates to the exact 10-K filing excerpt stored in GCS."
    )

    # ==========================================
    # SLIDE 8: Solution Overview — The 4 Technical Pillars
    # ==========================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide8)
    add_header(slide8, "Solution Overview — The 4 Technical Pillars", "Enterprise-Grade Financial Architecture on Google Cloud")

    pillars = [
        ("Pillar 1: Agentic Core (ADK)", "Google Agent Development Kit RootOrchestrator with Supervisor pattern & specialized Search Sub-Agent.", COLOR_GCP_BLUE),
        ("Pillar 2: Decoupled Math Engine", "Sandboxed Python execution solver (calculation_engine.py) for 100% computational accuracy.", COLOR_GCP_GREEN),
        ("Pillar 3: Hybrid Search RAG", "BigQuery SQL (structured metrics) + Vertex AI Search (unstructured 10-K vectors/BM25).", COLOR_GCP_YELLOW),
        ("Pillar 4: Security Perimeter", "VPC-SC perimeter, Model Armor ingress/egress callbacks, DLP PII redaction, & HITL gate.", COLOR_GCP_RED)
    ]

    for i, (title, desc, color) in enumerate(pillars):
        col = i % 2
        row = i // 2
        left = 0.8 + (col * 6.0)
        top = 1.8 + (row * 2.5)
        add_card(slide8, left, top, 5.7, 2.25, title, border_color=color)
        add_bullet_points(slide8, left + 0.2, top + 0.7, 5.3, 1.4, [desc], font_size=12)

    add_footer(slide8)
    set_speaker_notes(slide8,
        "Our architecture stands on four technical pillars: 1) ADK Supervisor for agentic orchestration, 2) Decoupled Python calculation engine for math, 3) Hybrid RAG combining BigQuery SQL and Vertex AI Search, and 4) Model Armor and VPC-SC for security."
    )

    # ==========================================
    # SLIDE 9: End-to-End System Architecture
    # ==========================================
    slide9 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide9)
    add_header(slide9, "End-to-End System Architecture Overview", "Detailed Data Flow across Presentation, API, ADK Core, & GCP Data Foundation")

    add_card(slide9, 0.8, 1.8, 11.733, 4.8, "System Flow (Refer to docs/architectural_diagrams.md Diagram 1)")
    add_bullet_points(slide9, 1.0, 2.4, 11.3, 4.0, [
        "Presentation Layer: React 18 / TypeScript Web UI executing Split-Pane context & A2UI JSON specification rendering.",
        "API & Session Layer: FastAPI on Cloud Run invoking AppController (app/app_controller.py) & Persistent Session Store.",
        "Security Ingress: Model Armor callback (model_armor_before_model_callback) inspecting prompt streams for jailbreaks & PII.",
        "Agentic Core: ADK RootOrchestrator (agent/root_orchestrator.py) applying System Constitution to route requests.",
        "Tools & Sub-Agents: Query BigQuery Tool, Sandboxed Calculation Tool, Search Sub-Agent, & HITL Export Tool.",
        "GCP Data Foundation: BigQuery Golden Tables, Vertex AI Search DataStore, & GCS Bucket Filings (gs://sec-analyst-sec-reports/)."
    ], font_size=12)

    add_footer(slide9)
    set_speaker_notes(slide9,
        "This end-to-end diagram illustrates the full execution path from the user prompt down to GCP storage layers. Model Armor inspects prompts at ingress before handing off to the ADK RootOrchestrator, which routes tasks to specialized tools and sub-agents."
    )

    # ==========================================
    # SLIDE 10: Layer Deep Dive 1 — Presentation Layer & A2UI
    # ==========================================
    slide10 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide10)
    add_header(slide10, "Layer Deep Dive 1 — Presentation Layer & A2UI Protocol", "React 18, Split-Pane Grounded Citations, & Structured Visual Protocols")

    add_card(slide10, 0.8, 1.8, 5.7, 4.8, "Presentation Architecture & Capabilities")
    add_bullet_points(slide10, 1.0, 2.5, 5.3, 3.8, [
        "React 18 + TypeScript: High-performance component architecture.",
        "Split-Pane Citation Engine: Real-time synchronization between narrative text and full SEC 10-K document viewer.",
        "Source Tracing: Links directly back to Ticker, Fiscal Year, Form 10-K, and Section (Item 7 MD&A).",
        "HITL Gate Modal: Requires explicit human sign-off before exporting generated reports."
    ], font_size=12)

    add_card(slide10, 6.8, 1.8, 5.7, 4.8, "Technical Trade-Off Matrix")
    add_card(slide10, 7.0, 2.4, 5.3, 1.8, "Selected Approach: A2UI Protocol", bg_color=COLOR_CARD_GREEN_BG, border_color=COLOR_GCP_GREEN)
    add_bullet_points(slide10, 7.1, 2.9, 5.1, 1.2, ["Emits structured A2UI JSON specs for dynamic tables & bar charts. Low latency, clean UI separation."], font_size=11)

    add_card(slide10, 7.0, 4.4, 5.3, 1.8, "Rejected Alternative: Raw Markdown Tables", bg_color=COLOR_CARD_RED_BG, border_color=COLOR_GCP_RED)
    add_bullet_points(slide10, 7.1, 4.9, 5.1, 1.2, ["Raw markdown tables pollute prompt context, break on long text, and cannot render dynamic charts."], font_size=11)

    add_footer(slide10)
    set_speaker_notes(slide10,
        "For Layer 1, we made a key architectural decision: instead of returning raw markdown table text from the LLM, we engineered the A2UI Protocol. The model outputs structured JSON specifications, which React renders into interactive financial tables and charts."
    )

    # ==========================================
    # SLIDE 11: Layer Deep Dive 2 — Agentic Core & ADK
    # ==========================================
    slide11 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide11)
    add_header(slide11, "Layer Deep Dive 2 — Agentic Core & ADK Orchestration", "Google Agent Development Kit (ADK) Supervisor Pattern & Intent Routing")

    add_card(slide11, 0.8, 1.8, 5.7, 4.8, "ADK Orchestration Logic")
    add_bullet_points(slide11, 1.0, 2.5, 5.3, 3.8, [
        "ADK RootOrchestrator: Built on Google ADK LlmAgent & Runner primitives (agent/root_orchestrator.py).",
        "Supervisor Pattern: Evaluates user intent and dynamically routes requests to tools and specialized sub-agents.",
        "Search Sub-Agent Decoupling: Offloads complex qualitative filing searches to agent/subagents/search_subagent.py.",
        "Persistent Session Storage: Maintains session state across multi-turn user flows without context loss."
    ], font_size=12)

    add_card(slide11, 6.8, 1.8, 5.7, 4.8, "Technical Trade-Off Matrix")
    add_card(slide11, 7.0, 2.4, 5.3, 1.8, "Selected Approach: ADK Supervisor Pattern", bg_color=COLOR_CARD_GREEN_BG, border_color=COLOR_GCP_GREEN)
    add_bullet_points(slide11, 7.1, 2.9, 5.1, 1.2, ["Modular tool routing, clear agent boundaries, specialized context isolation, & extensible tool registration."], font_size=11)

    add_card(slide11, 7.0, 4.4, 5.3, 1.8, "Rejected Alternative: Monolithic Script", bg_color=COLOR_CARD_RED_BG, border_color=COLOR_GCP_RED)
    add_bullet_points(slide11, 7.1, 4.9, 5.1, 1.2, ["Single monolithic prompts suffer from tool selection confusion, high context contamination, & brittle routing."], font_size=11)

    add_footer(slide11)
    set_speaker_notes(slide11,
        "In Layer 2, we chose the Google ADK Supervisor Pattern over a monolithic prompt. ADK's modular structure allows the RootOrchestrator to dynamically route metric queries to SQL, math to Python, and text searches to our Search Sub-Agent."
    )

    # ==========================================
    # SLIDE 12: Layer Deep Dive 3 — Hybrid Search RAG Layer
    # ==========================================
    slide12 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide12)
    add_header(slide12, "Layer Deep Dive 3 — Hybrid Search RAG Layer", "Dual-Path Retrieval: Merging BigQuery SQL Metrics & Vertex AI Search Vectors")

    add_card(slide12, 0.8, 1.8, 5.7, 4.8, "Dual-Path Retrieval Engine")
    add_bullet_points(slide12, 1.0, 2.5, 5.3, 3.8, [
        "Path 1 (Structured): BigQuery SQL tool directly queries exact numerical metrics from golden financial tables.",
        "Path 2 (Unstructured): Vertex AI Search DataStore performs hybrid search (dense semantic vectors + BM25 keyword matching) over 120+ SEC 10-K filings.",
        "Query Formulation: Strips conversational noise and anchors queries by Ticker and Fiscal Year.",
        "Parallel Highlighting: Orchestrates LLM calls to annotate exact source excerpts with <mark> tags."
    ], font_size=12)

    add_card(slide12, 6.8, 1.8, 5.7, 4.8, "Technical Trade-Off Matrix")
    add_card(slide12, 7.0, 2.4, 5.3, 1.8, "Selected: Hybrid (BigQuery SQL + Vertex AI Search)", bg_color=COLOR_CARD_GREEN_BG, border_color=COLOR_GCP_GREEN)
    add_bullet_points(slide12, 7.1, 2.9, 5.1, 1.2, ["Combines high-precision quantitative metrics from SQL with rich qualitative narrative context from 10-K vectors."], font_size=11)

    add_card(slide12, 7.0, 4.4, 5.3, 1.8, "Rejected Alternative: Pure Vector Retrieval", bg_color=COLOR_CARD_RED_BG, border_color=COLOR_GCP_RED)
    add_bullet_points(slide12, 7.1, 4.9, 5.1, 1.2, ["Pure vector search struggles with exact tabular metric lookups and risks retrieving outdated fiscal year rows."], font_size=11)

    add_footer(slide12)
    set_speaker_notes(slide12,
        "Our retrieval architecture uses a Dual-Path engine. We do NOT rely on vector search alone. Quantitative queries hit BigQuery SQL for exact numbers, while qualitative risk queries hit Vertex AI Hybrid Search, merging dense vectors with BM25 keywords."
    )

    # ==========================================
    # SLIDE 13: Layer Deep Dive 4 — Decoupled Math Engine
    # ==========================================
    slide13 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide13)
    add_header(slide13, "Layer Deep Dive 4 — Decoupled Deterministic Math Engine", "Sandboxed Python Code Execution for Zero-LLM Numerical Calculations")

    add_card(slide13, 0.8, 1.8, 5.7, 4.8, "Zero-LLM Math Approach & Safeguards")
    add_bullet_points(slide13, 1.0, 2.5, 5.3, 3.8, [
        "Sandboxed Execution: agent/tools/calculation_engine.py executes all variance, % delta, and margin calculations in Python.",
        "Zero LLM Arithmetic: LLM handles intent reasoning; deterministic code handles math execution.",
        "Division-by-Zero Traps: Automated error trapping and Inf/Null handling for sparse financial datasets.",
        "Negative Baseline Handling: Adjusted normalization logic for meaningful YoY variance analysis."
    ], font_size=12)

    add_card(slide13, 6.8, 1.8, 5.7, 4.8, "Technical Trade-Off Matrix")
    add_card(slide13, 7.0, 2.4, 5.3, 1.8, "Selected: Sandboxed Python Execution Engine", bg_color=COLOR_CARD_GREEN_BG, border_color=COLOR_GCP_GREEN)
    add_bullet_points(slide13, 7.1, 2.9, 5.1, 1.2, ["Achieves 100% computational accuracy, zero hallucinated digits, & complete mathematical auditability."], font_size=11)

    add_card(slide13, 7.0, 4.4, 5.3, 1.8, "Rejected Alternative: LLM Arithmetic Reasoning", bg_color=COLOR_CARD_RED_BG, border_color=COLOR_GCP_RED)
    add_bullet_points(slide13, 7.1, 4.9, 5.1, 1.2, ["Probabilistic LLM arithmetic introduces floating-point errors, incorrect signs, & hallucinated financial outputs."], font_size=11)

    add_footer(slide13)
    set_speaker_notes(slide13,
        "In finance, 99% math accuracy is a failing grade. We decoupled all math from the LLM into a sandboxed Python solver. The model extracts parameters, but Python executes the calculations, guaranteeing 100% arithmetic accuracy."
    )

    # ==========================================
    # SLIDE 14: Layer Deep Dive 5 — Security & Model Armor
    # ==========================================
    slide14 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide14)
    add_header(slide14, "Layer Deep Dive 5 — Security Perimeter & Model Armor Integration", "Multi-Layered Defense-in-Depth for Enterprise Compliance")

    add_card(slide14, 0.8, 1.8, 5.7, 4.8, "Security Infrastructure")
    add_bullet_points(slide14, 1.0, 2.5, 5.3, 3.8, [
        "VPC Service Controls (VPC-SC): Isolated security perimeter preventing data exfiltration to public networks.",
        "Model Armor Callbacks: Ingress and egress hooks inspecting prompts for adversarial jailbreaks and indirect injection.",
        "DLP PII Scrubber: Real-time redaction of sensitive personal identifiable information.",
        "Least Privilege IAM: Enforces sec-analyst-sa service account permissions with granular role bindings."
    ], font_size=12)

    add_card(slide14, 6.8, 1.8, 5.7, 4.8, "Technical Trade-Off Matrix")
    add_card(slide14, 7.0, 2.4, 5.3, 1.8, "Selected: VPC-SC + Model Armor + HITL Gate", bg_color=COLOR_CARD_GREEN_BG, border_color=COLOR_GCP_GREEN)
    add_bullet_points(slide14, 7.1, 2.9, 5.1, 1.2, ["Enterprise-grade defense-in-depth, proactive jailbreak blocking, & strict human sign-off on report exports."], font_size=11)

    add_card(slide14, 7.0, 4.4, 5.3, 1.8, "Rejected Alternative: Direct Public API Access", bg_color=COLOR_CARD_RED_BG, border_color=COLOR_GCP_RED)
    add_bullet_points(slide14, 7.1, 4.9, 5.1, 1.2, ["Public API endpoints expose systems to prompt injection, data exfiltration, & unmonitored model outputs."], font_size=11)

    add_footer(slide14)
    set_speaker_notes(slide14,
        "Layer 5 provides enterprise trust. We deploy inside a VPC Service Controls perimeter and integrate Model Armor callbacks on both ingress and egress. Adversarial jailbreaks are blocked before reaching Gemini."
    )

    # ==========================================
    # SLIDE 15: Evaluation Architecture — Dual-Track Auditing
    # ==========================================
    slide15 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide15)
    add_header(slide15, "Evaluation Architecture — Dual-Track Auditing Framework", "Combining Deterministic Math Verification with LLM-as-a-Judge")

    add_card(slide15, 0.8, 1.8, 5.7, 4.8, "Track 1: Math Assertions")
    add_bullet_points(slide15, 1.0, 2.5, 5.3, 3.8, [
        "Deterministic Verification: Automated verification of numerical outputs and percentage delta calculations.",
        "Range & Bound Constraints: Enforces non-negative constraints on absolute metrics.",
        "Cross-Statement Consistency: Validates that balance sheet totals match cash flow starting figures.",
        "Pass Rate: 100% Pass Rate across evaluation test suites."
    ], font_size=12)

    add_card(slide15, 6.8, 1.8, 5.7, 4.8, "Track 2: Gemini 2.5 Flash Judge", border_color=COLOR_GCP_BLUE)
    add_bullet_points(slide15, 7.0, 2.5, 5.3, 3.8, [
        "Semantic Quality Auditing: Evaluates narrative reasoning paths and qualitative completeness.",
        "Grounded Citation Exactness: Measures precision and recall of cited 10-K paragraph excerpts.",
        "Trajectory Quality: Evaluates step-by-step tool selection logic in ADK RootOrchestrator.",
        "Explainable Verdicts: Judge model generates structured JSON verdicts with reasoning notes."
    ], font_size=12)

    add_footer(slide15)
    set_speaker_notes(slide15,
        "We expanded our evaluation framework into two distinct tracks: Track 1 runs deterministic code assertions to verify math, while Track 2 uses Gemini 2.5 Flash as an LLM-as-a-Judge to evaluate citation precision and trajectory quality."
    )

    # ==========================================
    # SLIDE 16: ADK Evaluation Harness & Benchmark
    # ==========================================
    slide16 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide16)
    add_header(slide16, "ADK Evaluation Harness & Golden Dataset Benchmark", "Automated Pytest Suite Execution against Curated SEC Test Cases")

    add_card(slide16, 0.8, 1.8, 5.7, 4.8, "Eval Harness Architecture (run_adk_eval_parallel.py)")
    add_bullet_points(slide16, 1.0, 2.5, 5.3, 3.8, [
        "Golden Dataset: eval/golden_dataset.json containing 100+ curated financial Q&A ground-truth pairs.",
        "Automated Runner: eval/run_adk_eval_parallel.py executes parallel ADK agent trajectories.",
        "Metrics Extracted: Math accuracy, grounded recall, ROUGE-L F1, LLM faithfulness, and execution latency.",
        "Regression Lock: Automated pytest assertions block pull requests if evaluation scores drop."
    ], font_size=12)

    add_card(slide16, 6.8, 1.8, 5.7, 4.8, "Empirical Benchmark Results", border_color=COLOR_GCP_GREEN)
    add_bullet_points(slide16, 7.0, 2.5, 5.3, 3.8, [
        "Math Accuracy %: 100.0% (22/22 test cases passed)",
        "LLM Faithfulness: 1.0000 (Zero ungrounded claims)",
        "Answer Relevance: 1.0000 (Complete intent alignment)",
        "Execution Error Rate: 0.0% (Zero unhandled exceptions)",
        "Average E2E Latency: 1.31s (1,306.92 ms average latency)"
    ], font_size=12)

    add_footer(slide16)
    set_speaker_notes(slide16,
        "Our ADK Evaluation Harness runs automated benchmarks against a golden dataset. In our latest run of 22 comprehensive test cases, we achieved 100% math accuracy, 1.0 faithfulness, 0% execution errors, and an average latency of 1.31 seconds."
    )

    # ==========================================
    # SLIDE 17: Production Telemetry & Quality Control
    # ==========================================
    slide17 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide17)
    add_header(slide17, "Production Telemetry & Continuous Quality Control", "Non-Blocking BigQuery Event Sink & System Refinement Flywheel")

    add_card(slide17, 0.8, 1.8, 5.7, 4.8, "BigQuery Telemetry Sink (Implemented)")
    add_bullet_points(slide17, 1.0, 2.5, 5.3, 3.8, [
        "Asynchronous Streaming: BigQueryTelemetrySink (agent/observability/telemetry_sink.py) streams telemetry without user latency impact.",
        "Granular Event Logging: Captures trace ID, model name, input/output/cached tokens, latency, cost, and tool counts.",
        "Live Table: Sec_edgar_telemetry.telemetry_events stores full operational audit logs.",
        "Cost Attribution: Real-time calculation of token spend and context caching savings."
    ], font_size=12)

    add_card(slide17, 6.8, 1.8, 5.7, 4.8, "Refinement Flywheel & Continuous Sampling", border_color=COLOR_GCP_BLUE)
    add_bullet_points(slide17, 7.0, 2.5, 5.3, 3.8, [
        "Incident Detection: Runtime edge cases trigger system rule codification.",
        "System Constitution Lock: Codifies permanent behavior rules in agent/constitution.py.",
        "Pytest Assertion Lock: Locks new test assertions in eval/test_eval_harness.py to prevent regressions.",
        "Roadmap (5% Sampling): Extensible design to sample 5% of live traffic for automated drift auditing."
    ], font_size=12)

    add_footer(slide17)
    set_speaker_notes(slide17,
        "For observability, our BigQuery Telemetry Sink streams execution metrics asynchronously with zero user latency impact. Any edge cases detected feed into our Refinement Flywheel, where rules are permanently locked into the system constitution."
    )

    # ==========================================
    # SLIDE 18: Total Cost of Ownership & Tokenomics
    # ==========================================
    slide18 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide18)
    add_header(slide18, "Total Cost of Ownership (TCO) & Tokenomics (CAG)", "Analyzing Labor Efficiency & Context Caching Inference Optimization")

    add_card(slide18, 0.8, 1.8, 5.7, 4.8, "Empirical Tokenomics & CAG Discount")
    add_bullet_points(slide18, 1.0, 2.5, 5.3, 3.8, [
        "75% Input Cost Discount: Vertex AI Context Caching (CAG) reduces gemini-2.5-pro cached input to $0.3125/1M tokens (vs $1.25/1M standard).",
        "Live Savings Tracking: Tracked automatically by CostTracker (agent/observability/cost_tracker.py).",
        "Pre-fill Latency Reduction: Bypasses re-processing static system constitution & SEC schema instructions on every turn.",
        "Proven Latency SLA: Benchmark average execution latency of 1.31s (well under 3.0s SLA limit)."
    ], font_size=12)

    add_card(slide18, 6.8, 1.8, 5.7, 4.8, "Annual Cost Comparison", border_color=COLOR_GCP_GREEN)
    add_bullet_points(slide18, 7.0, 2.5, 5.3, 3.8, [
        "Manual Process Labor: $270,000 / year (3,000 hours at $90/hr analyst rate).",
        "Grounded AI Agent Labor: $10,000 / year (Gemini token usage + Cloud Run hosting).",
        "Net Annual Cost Reduction: 96% reduction in operational expense.",
        "Net Value Realized: $260,000 net annual savings per financial deployment."
    ], font_size=12)

    add_footer(slide18)
    set_speaker_notes(slide18,
        "Here are the tokenomics. Vertex AI Context Caching delivers an exact 75% input token discount on Gemini 2.5 Pro. Combined with our 1.31s execution latency, this reduces annual operating cost from $270,000 in manual labor down to $10,000 in AI cost."
    )

    # ==========================================
    # SLIDE 19: System Scalability & Production Readiness
    # ==========================================
    slide19 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide19)
    add_header(slide19, "System Scalability & Production Readiness", "High-Availability Deployment & Enterprise Infrastructure")

    add_card(slide19, 0.8, 1.8, 5.7, 4.8, "Compute & Container Infrastructure")
    add_bullet_points(slide19, 1.0, 2.5, 5.3, 3.8, [
        "Cloud Run Container Scaling: Serverless container deployment with dynamic autoscaling from 0 to 10 instances.",
        "Cost vs Cold Start Trade-Off: Configured min-instances=1 for zero cold-start latency on peak trading hours.",
        "IaC Provisioning: 100% Terraform Infrastructure as Code (main.tf) provisioning Cloud Run, GCS, IAM, & Secret Manager.",
        "Zero Hardcoded Environment: Dynamic secret mounting from GCP Secret Manager."
    ], font_size=12)

    add_card(slide19, 6.8, 1.8, 5.7, 4.8, "Data Layer & Privacy Governance")
    add_bullet_points(slide19, 7.0, 2.5, 5.3, 3.8, [
        "BigQuery Partitioning: Partitioned golden financial tables for high-concurrency analytical SQL queries.",
        "Data Privacy Isolation: All prompts and 10-K data remain strictly isolated inside the customer GCP project perimeter.",
        "Zero Model Training: Vertex AI guarantees customer prompt data is never used for foundation model re-training.",
        "Enterprise RBAC: Granular IAM role bindings for sec-analyst-sa service account."
    ], font_size=12)

    add_footer(slide19)
    set_speaker_notes(slide19,
        "Regarding scalability and security: our system is deployed on Cloud Run via 100% Terraform IaC with min-instances set to zero cold-start latency. BigQuery is partitioned for concurrency, and Vertex AI guarantees zero data exfiltration or model training on customer prompts."
    )

    # ==========================================
    # SLIDE 20: Panel Defense & Technical Trade-Off Summary
    # ==========================================
    slide20 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide20)
    add_header(slide20, "Panel Defense & Technical Trade-Off Matrix", "Strategic Justifications across Core Architectural Domains")

    table_shape = slide20.shapes.add_table(7, 4, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8))
    table = table_shape.table
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(3.2)
    table.columns[2].width = Inches(3.0)
    table.columns[3].width = Inches(3.333)

    headers = ["Domain", "Selected Production Approach", "Rejected Alternative", "Technical Justification"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_GCP_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_PRIMARY

    matrix_rows = [
        ("Orchestration", "ADK Supervisor + Sub-Agents", "Monolithic Script", "Modular routing & context isolation"),
        ("Math Logic", "Sandboxed Python Code Tool", "LLM Arithmetic", "100% math precision; zero hallucination"),
        ("Retrieval", "Hybrid (BigQuery SQL + Vector)", "Pure Vector Search", "Preserves structured table metric fidelity"),
        ("Security", "Model Armor + VPC-SC + HITL", "Standard Public API", "Defense-in-depth against data exfiltration"),
        ("Performance", "Context Caching (CAG)", "Uncached RAG Pipeline", "75% token discount & lower TTFT latency"),
        ("Evaluation", "ADK Harness + LLM-as-Judge", "Manual Prompt Testing", "Automated regression lock on golden dataset")
    ]

    for i, row in enumerate(matrix_rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = COLOR_TEXT_PRIMARY if j < 2 else (COLOR_TEXT_MUTED if j == 2 else COLOR_GCP_BLUE)

    add_footer(slide20)
    set_speaker_notes(slide20,
        "This defense matrix summarizes our core architectural choices. In every domain—orchestration, math, retrieval, security, performance, and evaluation—we selected deterministic, modular, and secure approaches over brittle monolithic alternatives."
    )

    # ==========================================
    # SLIDE 21: Business Impact & Strategic ROI
    # ==========================================
    slide21 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide21)
    add_header(slide21, "Business Impact & Strategic ROI Analysis", "Transforming Financial Analysis into a High-ROI Strategic Advantage")

    roi_cards = [
        ("26x Net ROI", "Financial Return", "26x Net Return on AI Investment"),
        ("$260,000", "Annual Net Savings", "$270k manual labor reduced to $10k AI cost"),
        ("2,900 Hours", "Reclaimed Capacity", "1.5 FTE senior analysts unlocked for alpha strategy"),
        ("< 2 Weeks", "Payback Period", "Rapid time-to-value realization per deployment")
    ]

    for i, (val, title, desc) in enumerate(roi_cards):
        left = 0.8 + (i * 3.0)
        add_card(slide21, left, 1.8, 2.7, 2.4, border_color=COLOR_GCP_GREEN)
        tx = slide21.shapes.add_textbox(Inches(left + 0.1), Inches(1.9), Inches(2.5), Inches(2.2))
        tf = tx.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.text = val
        p0.font.size = Pt(32)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_GCP_GREEN
        
        p1 = tf.add_paragraph()
        p1.text = title
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_PRIMARY

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    add_card(slide21, 0.8, 4.5, 11.733, 2.2, "Executive Takeaway", border_color=COLOR_GCP_BLUE)
    add_bullet_points(slide21, 1.0, 5.0, 11.3, 1.5, [
        "The SEC EDGAR Analyst Agent successfully replaces a $270,000 manual document processing bottleneck with a 100% mathematically precise, grounded agentic solution.",
        "By combining Google ADK orchestration, sandboxed Python math execution, and Vertex AI Context Caching, the platform delivers enterprise reliability with a payback period of under 2 weeks.",
        "Ready for immediate production deployment across Financial Services Industry (FSI) portfolios."
    ], font_size=12)

    add_footer(slide21)
    set_speaker_notes(slide21,
        "To conclude: our SEC EDGAR Analyst Agent delivers a 26x net ROI, $260,000 in net annual savings, and a payback period of under 2 weeks. More importantly, it reclaims 2,900 hours of senior analyst capacity while guaranteeing zero mathematical hallucinations. Thank you, and I am now ready for panel Q&A."
    )

    return prs


if __name__ == "__main__":
    prs = create_deck()
    output_path = "/Users/cvwang/Documents/gcp/sec-edgar-analyst/SEC_EDGAR_Analyst_FDE_Capstone_v2.pptx"
    prs.save(output_path)
    print(f"SUCCESS: Created updated 21-slide presentation at: {output_path}")
